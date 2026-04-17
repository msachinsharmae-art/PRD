"""Generate AI for BA Role PPT using template formatting."""
import copy
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from lxml import etree

TEMPLATE = "C:/Timesheet ESS/AI Use Case - Product.pptx"
OUTPUT = "C:/Timesheet ESS/AI for BA Role - Product Team.pptx"

prs = Presentation(TEMPLATE)

# ── Helper: duplicate a slide ──
def duplicate_slide(prs, template_slide):
    """Duplicate a slide by copying its XML."""
    slide_layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    # Remove default placeholders from new slide
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)
    # Copy all elements from template
    for shape in template_slide.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)
    return new_slide


def set_textbox_content(slide, textbox_name, paragraphs_data):
    """Replace content of a named textbox with new paragraphs.
    paragraphs_data: list of (text, bold, font_size_pt, is_bullet)
    """
    for shape in slide.shapes:
        if shape.name == textbox_name and shape.has_text_frame:
            tf = shape.text_frame
            # Clear existing paragraphs
            for _ in range(len(tf.paragraphs) - 1):
                p = tf.paragraphs[-1]._p
                p.getparent().remove(p)
            # Set first paragraph
            if paragraphs_data:
                first = paragraphs_data[0]
                p = tf.paragraphs[0]
                p.clear()
                run = p.add_run()
                run.text = first[0]
                run.font.bold = first[1]
                if first[2]:
                    run.font.size = Pt(first[2])
            # Add remaining paragraphs
            for i, (text, bold, size, is_bullet) in enumerate(paragraphs_data[1:]):
                p = tf.paragraphs[0]._p
                new_p = copy.deepcopy(p)
                # Clear runs in the copy
                nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                for r in new_p.findall('.//a:r', nsmap):
                    new_p.remove(r)
                # Add new run
                r_elem = etree.SubElement(new_p, '{http://schemas.openxmlformats.org/drawingml/2006/main}r')
                rPr = etree.SubElement(r_elem, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                rPr.set('lang', 'en-US')
                rPr.set('dirty', '0')
                if bold:
                    rPr.set('b', '1')
                if size:
                    rPr.set('sz', str(int(size * 100)))
                t_elem = etree.SubElement(r_elem, '{http://schemas.openxmlformats.org/drawingml/2006/main}t')
                t_elem.text = text
                tf.paragraphs[0]._p.getparent().append(new_p)
            return True
    return False


def update_title(slide, title_text):
    """Update the 'Text 2' shape (slide title like 'AI Use Case 1')."""
    for shape in slide.shapes:
        if shape.name == "Text 2" and shape.has_text_frame:
            tf = shape.text_frame
            # Clear all paragraphs first, then set text
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = ""
            # Set first paragraph
            if tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].text = title_text
            else:
                run = tf.paragraphs[0].add_run()
                run.text = title_text
            return


def update_slide_number(slide, num_str):
    """Update the slide number shape."""
    for shape in slide.shapes:
        if "128" in shape.name and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = num_str
            return


def get_content_textbox(slide):
    """Find the main content textbox name."""
    for shape in slide.shapes:
        if shape.name.startswith("TextBox") and shape.has_text_frame:
            if len(shape.text_frame.paragraphs) > 3:
                return shape.name
    return None


# ── Content definition ──
# Each slide: (title, slide_num, [(text, bold, font_size, is_bullet), ...])

slides_content = [
    # Slide 2: Why AI for BAs
    ("AI for BAs - Why Now?", "01", [
        ("Why AI for Business Analysts?", True, 18, False),
        ("", False, 12, False),
        ("The BA Challenge Today", True, 14, False),
        ("\u2022 Too much time on documentation, too little on analysis", False, 12, True),
        ("\u2022 Information scattered across tools, teams, and channels", False, 12, True),
        ("\u2022 Repetitive tasks \u2014 ticket writing, status reports, meeting notes", False, 12, True),
        ("\u2022 Stakeholder communication takes multiple drafts", False, 12, True),
        ("", False, 12, False),
        ("The AI Opportunity", True, 14, False),
        ("\u2022 Automate the repetitive, focus on the strategic", False, 12, True),
        ("\u2022 Reduce first-draft time by 50-70%", False, 12, True),
        ("\u2022 Improve consistency and coverage in requirements", False, 12, True),
        ("\u2022 Free up time for stakeholder engagement & product thinking", False, 12, True),
    ]),

    # Slide 3: AI Across BA Workflow
    ("AI Across the BA Workflow", "02", [
        ("AI Across the BA Workflow", True, 18, False),
        ("Discovery \u2192 Definition \u2192 Delivery \u2192 Review", True, 13, False),
        ("", False, 12, False),
        ("Discovery Phase", True, 14, False),
        ("\u2022 Without AI: Manual research, hours of reading", False, 12, True),
        ("\u2022 With AI: AI summarizes reports, competitors, regulations", False, 12, True),
        ("", False, 12, False),
        ("Definition Phase", True, 14, False),
        ("\u2022 Without AI: Blank-page PRDs, missed edge cases", False, 12, True),
        ("\u2022 With AI: AI drafts PRDs, flags gaps automatically", False, 12, True),
        ("", False, 12, False),
        ("Delivery Phase", True, 14, False),
        ("\u2022 Without AI: Manual ticket writing, back-and-forth", False, 12, True),
        ("\u2022 With AI: AI generates user stories & test cases", False, 12, True),
        ("", False, 12, False),
        ("Review Phase", True, 14, False),
        ("\u2022 Without AI: Manual data crunching for insights", False, 12, True),
        ("\u2022 With AI: AI analyzes data, highlights trends", False, 12, True),
    ]),

    # Slide 4: Use Case 1 - PRD
    ("Use Case 1", "03", [
        ("PRD & Requirements Creation", True, 18, False),
        ("", False, 12, False),
        ("Problem", True, 14, False),
        ("Writing a PRD from scratch takes hours and often misses edge cases.", False, 12, False),
        ("", False, 12, False),
        ("AI Solution", True, 14, False),
        ("\u2022 Input: rough notes, stakeholder call points, feature brief", False, 12, True),
        ("\u2022 Output: structured PRD with objective, scope, user stories, acceptance criteria, risks", False, 12, True),
        ("", False, 12, False),
        ("Example Prompt", True, 14, False),
        ("\"I need a PRD for an employee leave management module. Employees should apply for leave, managers approve/reject, HR views reports. Include edge cases.\"", False, 11, False),
        ("", False, 12, False),
        ("Impact", True, 14, False),
        ("PRD first draft in 10 minutes instead of 2-3 hours", False, 12, False),
    ]),

    # Slide 5: Use Case 2 - User Stories
    ("Use Case 2", "04", [
        ("User Story & Acceptance Criteria", True, 18, False),
        ("", False, 12, False),
        ("Problem", True, 14, False),
        ("Writing detailed user stories with all edge cases is time-consuming.", False, 12, False),
        ("", False, 12, False),
        ("AI Solution", True, 14, False),
        ("\u2022 Generates stories in standard format (As a... I want... So that...)", False, 12, True),
        ("\u2022 Adds acceptance criteria in Given/When/Then", False, 12, True),
        ("\u2022 Suggests edge cases you might miss", False, 12, True),
        ("", False, 12, False),
        ("Example Output", True, 14, False),
        ("Story: As a manager, I want to approve/reject leave requests so that team availability is maintained", False, 11, False),
        ("AC1: Given a pending leave request, When manager clicks Approve, Then leave balance is deducted and employee is notified", False, 11, False),
        ("AC2: Given a pending request with zero balance, When manager views it, Then a warning is displayed", False, 11, False),
        ("", False, 12, False),
        ("Impact: 5x faster story writing with better coverage", True, 12, False),
    ]),

    # Slide 6: Use Case 3 - Research
    ("Use Case 3", "05", [
        ("Competitive & Market Research", True, 18, False),
        ("", False, 12, False),
        ("Problem", True, 14, False),
        ("Researching competitors and market trends takes days.", False, 12, False),
        ("", False, 12, False),
        ("AI Solution", True, 14, False),
        ("\u2022 Summarize competitor product pages and features", False, 12, True),
        ("\u2022 Compare feature sets across multiple products", False, 12, True),
        ("\u2022 Explain regulatory/compliance requirements in simple language", False, 12, True),
        ("\u2022 Identify market gaps and opportunities", False, 12, True),
        ("", False, 12, False),
        ("Example", True, 14, False),
        ("\"Compare the leave management features of Zimyo, Keka, and GreytHR. Create a feature matrix.\"", False, 11, False),
        ("", False, 12, False),
        ("Impact", True, 14, False),
        ("Research that took 2 days now takes 2 hours", False, 12, False),
    ]),

    # Slide 7: Use Case 4 - Process Mapping
    ("Use Case 4", "06", [
        ("Process Mapping & Documentation", True, 18, False),
        ("", False, 12, False),
        ("Problem", True, 14, False),
        ("Creating flowcharts and SOPs is tedious and often delayed.", False, 12, False),
        ("", False, 12, False),
        ("AI Solution", True, 14, False),
        ("\u2022 Describe a process in plain English \u2192 get a flowchart (Mermaid/PlantUML)", False, 12, True),
        ("\u2022 Convert informal knowledge into structured SOPs", False, 12, True),
        ("\u2022 Generate BPMN-style process descriptions", False, 12, True),
        ("", False, 12, False),
        ("Example Prompt", True, 14, False),
        ("\"Create a flowchart for employee reimbursement: Employee submits claim \u2192 Manager approves \u2192 Finance verifies \u2192 Payment processed. Include rejection paths.\"", False, 11, False),
        ("", False, 12, False),
        ("Impact", True, 14, False),
        ("Process documentation created in minutes, not days", False, 12, False),
    ]),

    # Slide 8: Use Case 5 - Meeting Notes
    ("Use Case 5", "07", [
        ("Meeting Notes & Action Items", True, 18, False),
        ("", False, 12, False),
        ("Problem", True, 14, False),
        ("Note-taking during meetings is distracting, action items get lost.", False, 12, False),
        ("", False, 12, False),
        ("AI Solution", True, 14, False),
        ("\u2022 Paste meeting transcript or rough notes", False, 12, True),
        ("\u2022 AI structures into: Key Decisions, Action Items (with owners), Open Questions", False, 12, True),
        ("\u2022 Generate follow-up email draft", False, 12, True),
        ("", False, 12, False),
        ("Example Output", True, 14, False),
        ("Key Decision: Sprint 14 will focus on payroll compliance module", False, 11, False),
        ("Action Item: Rahul to share PF calculation logic by March 18", False, 11, False),
        ("Open Question: Do we support multi-state PT in Phase 1?", False, 11, False),
        ("", False, 12, False),
        ("Impact: 15 minutes of post-meeting work reduced to 2 minutes", True, 12, False),
    ]),

    # Slide 9: Use Case 6 - Test Cases
    ("Use Case 6", "08", [
        ("Test Case Generation", True, 18, False),
        ("", False, 12, False),
        ("Problem", True, 14, False),
        ("BAs often support QA with test scenarios but coverage gaps remain.", False, 12, False),
        ("", False, 12, False),
        ("AI Solution", True, 14, False),
        ("\u2022 Feed acceptance criteria \u2192 get positive, negative, and edge case test scenarios", False, 12, True),
        ("\u2022 Covers boundary values, error states, permission checks", False, 12, True),
        ("\u2022 Structured format ready for test management tools", False, 12, True),
        ("", False, 12, False),
        ("Example", True, 14, False),
        ("Input: \"Employee can apply for leave only if balance > 0 and notice period is 2 days\"", False, 11, False),
        ("\u2022 TC1: Apply with sufficient balance and notice \u2192 Success", False, 11, True),
        ("\u2022 TC2: Apply with zero balance \u2192 Error shown", False, 11, True),
        ("\u2022 TC3: Apply with 1-day notice \u2192 Error shown", False, 11, True),
        ("\u2022 TC4: Apply on weekend/holiday \u2192 Behavior check", False, 11, True),
        ("", False, 12, False),
        ("Impact: 3x more test coverage, 60% less effort", True, 12, False),
    ]),

    # Slide 10: Use Case 7 - Stakeholder Communication
    ("Use Case 7", "09", [
        ("Stakeholder Communication", True, 18, False),
        ("", False, 12, False),
        ("AI Solution", True, 14, False),
        ("", False, 12, False),
        ("Sprint Release Notes", True, 13, False),
        ("\u2022 Generate user-friendly notes from Jira tickets", False, 12, True),
        ("", False, 12, False),
        ("Escalation Emails", True, 13, False),
        ("\u2022 Professional, structured escalation drafts", False, 12, True),
        ("", False, 12, False),
        ("Status Reports", True, 13, False),
        ("\u2022 Paste raw updates \u2192 get formatted weekly report", False, 12, True),
        ("", False, 12, False),
        ("Technical to Business Translation", True, 13, False),
        ("\u2022 Explain API limitations to business OR translate business needs into dev specs", False, 12, True),
        ("", False, 12, False),
        ("Impact: Communication quality up, turnaround time down", True, 12, False),
    ]),

    # Slide 11: Use Case 8 - Data Analysis
    ("Use Case 8", "10", [
        ("Data Analysis & Insights", True, 18, False),
        ("", False, 12, False),
        ("Problem", True, 14, False),
        ("BAs need to analyze data but writing queries and interpreting results is slow.", False, 12, False),
        ("", False, 12, False),
        ("AI Solution", True, 14, False),
        ("\u2022 Write SQL queries from plain English questions", False, 12, True),
        ("\u2022 Create Excel formulas for complex calculations", False, 12, True),
        ("\u2022 Analyze support tickets/feedback to find patterns", False, 12, True),
        ("\u2022 Suggest KPIs for new features", False, 12, True),
        ("", False, 12, False),
        ("Example", True, 14, False),
        ("\"Write a SQL query to find employees who applied for more than 3 leaves in the last 30 days, grouped by department\"", False, 11, False),
        ("", False, 12, False),
        ("Impact", True, 14, False),
        ("Data tasks that took hours now take minutes", False, 12, False),
    ]),

    # Slide 12: Day-to-Day Automation
    ("Day-to-Day Automation", "11", [
        ("Day-to-Day Task Automation", True, 18, False),
        ("", False, 12, False),
        ("\u2022 Jira ticket writing \u2192 Generate from rough notes \u2192 60-70% time saved", False, 12, True),
        ("\u2022 PRD drafting \u2192 Structured first draft \u2192 50-60% time saved", False, 12, True),
        ("\u2022 Meeting follow-ups \u2192 Auto-structured notes & emails \u2192 70% time saved", False, 12, True),
        ("\u2022 Sprint planning prep \u2192 Backlog summary & priority suggestions \u2192 30-40% time saved", False, 12, True),
        ("\u2022 Bug triage \u2192 Impact analysis & root cause hints \u2192 40% time saved", False, 12, True),
        ("\u2022 Status reporting \u2192 Auto-generate from updates \u2192 50% time saved", False, 12, True),
        ("\u2022 Test case support \u2192 Generate from acceptance criteria \u2192 60% time saved", False, 12, True),
        ("\u2022 Data queries \u2192 SQL/Excel from plain English \u2192 80% time saved", False, 12, True),
    ]),

    # Slide 13: Before vs After
    ("Before vs After", "12", [
        ("Before vs After AI Adoption", True, 18, False),
        ("", False, 12, False),
        ("\u2022 PRD creation time: 4-6 hours \u2192 1-2 hours", False, 12, True),
        ("\u2022 User stories per sprint: 15-20 \u2192 30-40", False, 12, True),
        ("\u2022 Edge case coverage: 60-70% \u2192 85-95%", False, 12, True),
        ("\u2022 Documentation backlog: Always behind \u2192 Up to date", False, 12, True),
        ("\u2022 Research turnaround: 2-3 days \u2192 3-4 hours", False, 12, True),
        ("\u2022 Stakeholder rework cycles: 3-4 rounds \u2192 1-2 rounds", False, 12, True),
    ]),

    # Slide 14: Tools
    ("Recommended Tools", "13", [
        ("Recommended Tools", True, 18, False),
        ("", False, 12, False),
        ("Claude (Anthropic)", True, 13, False),
        ("\u2022 Best for: PRDs, analysis, long documents, coding | Free / Pro", False, 12, True),
        ("", False, 12, False),
        ("ChatGPT (OpenAI)", True, 13, False),
        ("\u2022 Best for: General tasks, quick drafts | Free / Plus", False, 12, True),
        ("", False, 12, False),
        ("Notion AI", True, 13, False),
        ("\u2022 Best for: In-tool docs, wikis, notes | Add-on", False, 12, True),
        ("", False, 12, False),
        ("GitHub Copilot / Claude Code", True, 13, False),
        ("\u2022 Best for: Technical specs, API understanding | Subscription", False, 12, True),
        ("", False, 12, False),
        ("Otter.ai / Fireflies", True, 13, False),
        ("\u2022 Best for: Meeting transcription | Free / Paid", False, 12, True),
        ("", False, 12, False),
        ("Gamma.app", True, 13, False),
        ("\u2022 Best for: AI-generated presentations | Free / Paid", False, 12, True),
    ]),

    # Slide 15: Best Practices
    ("Best Practices", "14", [
        ("Best Practices", True, 18, False),
        ("", False, 12, False),
        ("1. AI = First Draft, Not Final Answer", True, 13, False),
        ("    Always review and refine AI output before sharing", False, 12, False),
        ("", False, 12, False),
        ("2. Give Context", True, 13, False),
        ("    The more specific your input, the better the output", False, 12, False),
        ("", False, 12, False),
        ("3. Build Prompt Templates", True, 13, False),
        ("    Reuse prompts for recurring tasks (PRD, stories, test cases)", False, 12, False),
        ("", False, 12, False),
        ("4. Don't Share Sensitive Data", True, 13, False),
        ("    Avoid pasting PII, credentials, or confidential business data", False, 12, False),
        ("", False, 12, False),
        ("5. Start Small", True, 13, False),
        ("    Pick 2-3 daily tasks, build the habit, then expand", False, 12, False),
        ("", False, 12, False),
        ("6. Document What Works", True, 13, False),
        ("    Maintain a team prompt library", False, 12, False),
    ]),

    # Slide 16: 30-Day Action Plan
    ("30-Day Action Plan", "15", [
        ("Getting Started \u2014 30-Day Action Plan", True, 18, False),
        ("", False, 12, False),
        ("Week 1", True, 14, False),
        ("Pick 2 repetitive tasks \u2192 start using AI for first drafts", False, 12, False),
        ("", False, 12, False),
        ("Week 2", True, 14, False),
        ("Build 3-5 reusable prompt templates for your common tasks", False, 12, False),
        ("", False, 12, False),
        ("Week 3", True, 14, False),
        ("Share results with team \u2192 collect feedback \u2192 refine prompts", False, 12, False),
        ("", False, 12, False),
        ("Week 4", True, 14, False),
        ("Identify next set of tasks to automate \u2192 create team prompt library", False, 12, False),
        ("", False, 14, False),
        ("Goal: Save 5-8 hours per week within 30 days", True, 16, False),
    ]),

    # Slide 17: Q&A
    ("Q&A", "16", [
        ("Questions & Discussion", True, 18, False),
        ("", False, 14, False),
        ("", False, 14, False),
        ("Key Takeaway", True, 16, False),
        ("", False, 12, False),
        ("AI doesn't replace the BA \u2014 it removes the grunt work so you can focus on what matters:", False, 14, False),
        ("", False, 12, False),
        ("Understanding users", True, 14, False),
        ("Shaping products", True, 14, False),
        ("Driving decisions", True, 14, False),
    ]),
]

# ── Reference slide (slide index 1 = Slide 2) for cloning ──
ref_slide = prs.slides[1]

# ── Step 1: Update title slide tagline ──
title_slide = prs.slides[0]
for shape in title_slide.shapes:
    if shape.name == "Google Shape;98;p1" and shape.has_text_frame:
        for run in shape.text_frame.paragraphs[0].runs:
            run.text = "AI-Powered Business Analysis for Product Teams"

# ── Step 2: Update existing content slides (indices 1-7) and add new ones ──
total_existing_content = 7  # slides index 1-7
total_needed = len(slides_content)

# Update existing slides (up to 7)
for i in range(min(total_existing_content, total_needed)):
    slide = prs.slides[i + 1]  # offset by 1 for title slide
    title, num, content = slides_content[i]

    # Update title
    update_title(slide, title)

    # Update slide number
    update_slide_number(slide, num)

    # Find and update content textbox
    content_tb = get_content_textbox(slide)
    if content_tb:
        set_textbox_content(slide, content_tb, content)

    # Clear ALL secondary textboxes and extra Text 2 shapes
    found_text2 = False
    for shape in slide.shapes:
        if shape.name == "Text 2" and shape.has_text_frame:
            if found_text2:
                # Extra Text 2 — clear it
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""
            found_text2 = True
        elif shape.name.startswith("TextBox") and shape.has_text_frame and shape.name != content_tb:
            if len(shape.text_frame.paragraphs) > 1:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""

# Add additional slides by duplicating reference slide
for i in range(total_existing_content, total_needed):
    title, num, content = slides_content[i]
    new_slide = duplicate_slide(prs, ref_slide)

    update_title(new_slide, title)
    update_slide_number(new_slide, num)

    content_tb = get_content_textbox(new_slide)
    if content_tb:
        set_textbox_content(new_slide, content_tb, content)

# ── Step 3: Move closing slide to the end ──
# The closing slide is at index 8 (original slide 9).
# New slides were appended after it, so we need to reorder.
# In python-pptx, we manipulate the XML to reorder.
pres_elem = prs.part._element
nsmap_p = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
sldIdLst = pres_elem.find('.//p:sldIdLst', nsmap_p)
sldIds = list(sldIdLst)

# Original order: [title, s2..s8, closing, new_s9..new_s17]
# We want: [title, s2..s8, new_s9..new_s17, closing]
# closing is at index 8
if len(sldIds) > 9:
    closing_sldId = sldIds[8]  # the closing slide
    sldIdLst.remove(closing_sldId)
    sldIdLst.append(closing_sldId)

prs.save(OUTPUT)
print(f"PPT saved to: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
