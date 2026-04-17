"""Add flowchart visuals to right side of each content slide."""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INPUT = "C:/Timesheet ESS/AI for BA Role - Product Team.pptx"
OUTPUT = "C:/Timesheet ESS/AI for BA Role - Product Team.pptx"

prs = Presentation(INPUT)

# ── Color palette (from template) ──
DARK_TEAL = RGBColor(0x03, 0x19, 0x2E)
TEAL = RGBColor(0x00, 0x7B, 0x83)
LIGHT_TEAL = RGBColor(0xE0, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xFF, 0x8F, 0x00)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
BLUE = RGBColor(0x15, 0x65, 0xC0)
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
RED_ACCENT = RGBColor(0xE5, 0x39, 0x35)

# ── Right-side area config ──
# For wide-textbox slides: shrink textbox, place visual on right
# For two-column slides: use right column area
RIGHT_X = Inches(7.8)
RIGHT_Y = Inches(1.8)
RIGHT_W = Inches(4.8)
RIGHT_H = Inches(5.0)

def add_rounded_box(slide, left, top, width, height, text, fill_color, font_color=WHITE, font_size=11, bold=True):
    """Add a rounded rectangle with text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    # Adjust corner rounding
    shape.adjustments[0] = 0.15

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = "Poppins"
    tf.auto_size = None
    shape.text_frame.margin_left = Pt(4)
    shape.text_frame.margin_right = Pt(4)
    shape.text_frame.margin_top = Pt(4)
    shape.text_frame.margin_bottom = Pt(4)
    shape.text_frame.word_wrap = True
    return shape


def add_arrow_down(slide, cx, top, length, color=TEAL):
    """Add a downward arrow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, cx - Inches(0.2), top, Inches(0.4), length
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_arrow_right(slide, left, cy, length, color=TEAL):
    """Add a right arrow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, cy - Inches(0.15), length, Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_circle(slide, left, top, size, text, fill_color, font_color=WHITE, font_size=10):
    """Add a circle with text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = True
    run.font.name = "Poppins"
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    return shape


def add_text_label(slide, left, top, width, height, text, font_size=10, color=DARK_TEAL, bold=False, align=PP_ALIGN.CENTER):
    """Add a text label."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Poppins"
    return txBox


def shrink_wide_textbox(slide, max_width_inches=7.2):
    """Shrink wide textboxes to make room for visuals."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name.startswith("TextBox"):
            if Emu(shape.width).inches > 9:
                shape.width = Inches(max_width_inches)


def clear_right_textbox(slide):
    """Clear or remove secondary right-side textboxes."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name.startswith("TextBox"):
            if Emu(shape.left).inches > 6:
                # Clear it
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""


# ═══════════════════════════════════════════════
# SLIDE 2: Why AI for BAs — Time comparison visual
# ═══════════════════════════════════════════════
slide = prs.slides[1]
shrink_wide_textbox(slide)

bx = Inches(8.0)
by = Inches(2.0)

add_text_label(slide, bx, by, Inches(4.5), Inches(0.5), "BA Workflow Transformation", 14, DARK_TEAL, True)

# Before column
add_text_label(slide, bx, by + Inches(0.7), Inches(2.0), Inches(0.4), "BEFORE AI", 11, RED_ACCENT, True)
bars_before = [("Docs", 3.5), ("Research", 3.0), ("Analysis", 1.5), ("Strategy", 0.8)]
for i, (label, width) in enumerate(bars_before):
    y = by + Inches(1.2 + i * 0.55)
    add_text_label(slide, bx, y, Inches(1.0), Inches(0.35), label, 8, DARK_TEAL, False, PP_ALIGN.RIGHT)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx + Inches(1.1), y, Inches(width), Inches(0.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xEF, 0x9A, 0x9A)
    bar.line.fill.background()
    bar.shadow.inherit = False

# After column
add_text_label(slide, bx, by + Inches(3.5), Inches(2.0), Inches(0.4), "AFTER AI", 11, GREEN, True)
bars_after = [("Docs", 1.2), ("Research", 1.0), ("Analysis", 2.5), ("Strategy", 3.2)]
for i, (label, width) in enumerate(bars_after):
    y = by + Inches(4.0 + i * 0.55)
    add_text_label(slide, bx, y, Inches(1.0), Inches(0.35), label, 8, DARK_TEAL, False, PP_ALIGN.RIGHT)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx + Inches(1.1), y, Inches(width), Inches(0.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xA5, 0xD6, 0xA7)
    bar.line.fill.background()
    bar.shadow.inherit = False


# ═══════════════════════════════════════════════
# SLIDE 3: BA Workflow — 4-phase flow
# ═══════════════════════════════════════════════
slide = prs.slides[2]
shrink_wide_textbox(slide)

bx = Inches(8.2)
by = Inches(2.0)
box_w = Inches(3.8)
box_h = Inches(0.65)
gap = Inches(0.45)

phases = [
    ("Discovery", "AI summarizes & researches", TEAL),
    ("Definition", "AI drafts PRDs & stories", BLUE),
    ("Delivery", "AI generates tickets & tests", PURPLE),
    ("Review", "AI analyzes data & trends", GREEN),
]

for i, (phase, desc, color) in enumerate(phases):
    y = by + i * (box_h + gap + Inches(0.3))
    add_rounded_box(slide, bx, y, box_w, box_h, phase, color, WHITE, 13, True)
    add_text_label(slide, bx, y + box_h + Inches(0.02), box_w, Inches(0.3), desc, 9, DARK_TEAL, False)
    if i < len(phases) - 1:
        add_arrow_down(slide, bx + box_w / 2, y + box_h + Inches(0.25), Inches(0.35), color)


# ═══════════════════════════════════════════════
# SLIDE 4: PRD Creation — Flow diagram
# ═══════════════════════════════════════════════
slide = prs.slides[3]
clear_right_textbox(slide)

bx = Inches(7.0)
by = Inches(2.2)

steps = [
    ("Rough Notes\n& Ideas", LIGHT_TEAL, DARK_TEAL),
    ("AI Engine", TEAL, WHITE),
    ("Structured\nPRD", GREEN, WHITE),
]

box_w = Inches(2.8)
box_h = Inches(1.0)

for i, (text, fill, font_c) in enumerate(steps):
    y = by + i * (box_h + Inches(0.6))
    add_rounded_box(slide, bx + Inches(1.0), y, box_w, box_h, text, fill, font_c, 12, True)
    if i < len(steps) - 1:
        add_arrow_down(slide, bx + Inches(1.0) + box_w / 2, y + box_h + Inches(0.05), Inches(0.45), TEAL)

add_text_label(slide, bx + Inches(0.5), by + Inches(5.0), Inches(4.0), Inches(0.4),
               "10 min vs 2-3 hours", 13, ORANGE, True)


# ═══════════════════════════════════════════════
# SLIDE 5: User Stories — Input/Output flow
# ═══════════════════════════════════════════════
slide = prs.slides[4]
clear_right_textbox(slide)

bx = Inches(7.0)
by = Inches(2.0)

add_rounded_box(slide, bx + Inches(0.8), by, Inches(3.2), Inches(0.8), "Feature Brief", LIGHT_TEAL, DARK_TEAL, 12)
add_arrow_down(slide, bx + Inches(2.4), by + Inches(0.85), Inches(0.4), TEAL)

add_rounded_box(slide, bx + Inches(0.8), by + Inches(1.4), Inches(3.2), Inches(0.8), "AI Processing", TEAL, WHITE, 12)
add_arrow_down(slide, bx + Inches(2.4), by + Inches(2.25), Inches(0.4), TEAL)

# Output boxes
outputs = ["User Stories", "Acceptance Criteria", "Edge Cases"]
colors = [GREEN, BLUE, PURPLE]
for i, (txt, col) in enumerate(zip(outputs, colors)):
    add_rounded_box(slide, bx + Inches(0.2) + Inches(i * 1.6), by + Inches(2.9),
                    Inches(1.5), Inches(0.7), txt, col, WHITE, 9)

add_text_label(slide, bx + Inches(0.5), by + Inches(4.0), Inches(4.0), Inches(0.4),
               "5x faster with better coverage", 12, ORANGE, True)


# ═══════════════════════════════════════════════
# SLIDE 6: Research — Multiple sources funnel
# ═══════════════════════════════════════════════
slide = prs.slides[5]
clear_right_textbox(slide)

bx = Inches(7.2)
by = Inches(1.9)

sources = ["Competitors", "Market Data", "Regulations", "User Feedback"]
src_colors = [TEAL, BLUE, PURPLE, GREEN]
for i, (src, col) in enumerate(zip(sources, src_colors)):
    add_rounded_box(slide, bx + Inches(i % 2 * 2.2), by + Inches(int(i / 2) * 1.0),
                    Inches(2.0), Inches(0.65), src, col, WHITE, 10)

# Funnel arrow
add_arrow_down(slide, bx + Inches(2.1), by + Inches(2.1), Inches(0.5), TEAL)

# AI box
add_rounded_box(slide, bx + Inches(0.6), by + Inches(2.8), Inches(3.0), Inches(0.8), "AI Analysis", DARK_TEAL, WHITE, 13)
add_arrow_down(slide, bx + Inches(2.1), by + Inches(3.7), Inches(0.5), TEAL)

# Output
outputs_r = ["Feature Matrix", "Gap Analysis", "Insights"]
for i, txt in enumerate(outputs_r):
    add_rounded_box(slide, bx + Inches(i * 1.5), by + Inches(4.4),
                    Inches(1.4), Inches(0.6), txt, GREEN, WHITE, 9)


# ═══════════════════════════════════════════════
# SLIDE 7: Process Mapping — Text to Flowchart
# ═══════════════════════════════════════════════
slide = prs.slides[6]
clear_right_textbox(slide)

bx = Inches(7.2)
by = Inches(2.0)

add_rounded_box(slide, bx + Inches(0.5), by, Inches(3.5), Inches(0.7),
                "Plain English Description", LIGHT_TEAL, DARK_TEAL, 11)
add_arrow_down(slide, bx + Inches(2.25), by + Inches(0.75), Inches(0.4), TEAL)

add_rounded_box(slide, bx + Inches(0.5), by + Inches(1.3), Inches(3.5), Inches(0.7),
                "AI Process Engine", TEAL, WHITE, 12)
add_arrow_down(slide, bx + Inches(2.25), by + Inches(2.05), Inches(0.4), TEAL)

# Mini flowchart output
fc_y = by + Inches(2.7)
fc_boxes = ["Start", "Step 1", "Step 2", "End"]
fc_colors = [GREEN, BLUE, BLUE, GREEN]
for i, (txt, col) in enumerate(zip(fc_boxes, fc_colors)):
    x = bx + Inches(0.2) + Inches(i * 1.15)
    add_rounded_box(slide, x, fc_y, Inches(0.95), Inches(0.5), txt, col, WHITE, 8)
    if i < len(fc_boxes) - 1:
        add_arrow_right(slide, x + Inches(0.95), fc_y + Inches(0.25), Inches(0.2), TEAL)

add_text_label(slide, bx + Inches(0.5), fc_y + Inches(0.8), Inches(3.5), Inches(0.4),
               "Flowchart + SOP + BPMN", 11, ORANGE, True)


# ═══════════════════════════════════════════════
# SLIDE 8: Meeting Notes — Transcript to Actions
# ═══════════════════════════════════════════════
slide = prs.slides[7]
clear_right_textbox(slide)

bx = Inches(7.0)
by = Inches(2.0)

add_rounded_box(slide, bx + Inches(0.8), by, Inches(3.2), Inches(0.8),
                "Meeting Transcript\n/ Rough Notes", LIGHT_TEAL, DARK_TEAL, 11)
add_arrow_down(slide, bx + Inches(2.4), by + Inches(0.85), Inches(0.4), TEAL)

add_rounded_box(slide, bx + Inches(0.8), by + Inches(1.4), Inches(3.2), Inches(0.7),
                "AI Structuring", TEAL, WHITE, 12)
add_arrow_down(slide, bx + Inches(2.4), by + Inches(2.15), Inches(0.4), TEAL)

outputs_m = [
    ("Key Decisions", GREEN),
    ("Action Items", BLUE),
    ("Open Questions", PURPLE),
    ("Follow-up Email", ORANGE),
]
for i, (txt, col) in enumerate(outputs_m):
    add_rounded_box(slide, bx + Inches(0.3) + Inches((i % 2) * 2.3), by + Inches(2.8 + int(i / 2) * 0.9),
                    Inches(2.1), Inches(0.65), txt, col, WHITE, 10)

add_text_label(slide, bx + Inches(0.5), by + Inches(4.8), Inches(4.0), Inches(0.4),
               "15 min → 2 min", 13, ORANGE, True)


# ═══════════════════════════════════════════════
# SLIDE 9: Test Cases — Criteria to Test flow
# ═══════════════════════════════════════════════
slide = prs.slides[8]
shrink_wide_textbox(slide)

bx = Inches(8.0)
by = Inches(2.0)

add_rounded_box(slide, bx + Inches(0.5), by, Inches(3.5), Inches(0.7),
                "Acceptance Criteria", LIGHT_TEAL, DARK_TEAL, 11)
add_arrow_down(slide, bx + Inches(2.25), by + Inches(0.75), Inches(0.35), TEAL)

add_rounded_box(slide, bx + Inches(0.5), by + Inches(1.2), Inches(3.5), Inches(0.7),
                "AI Test Generator", TEAL, WHITE, 12)
add_arrow_down(slide, bx + Inches(2.25), by + Inches(1.95), Inches(0.35), TEAL)

tc_types = [
    ("Positive\nTests", GREEN),
    ("Negative\nTests", RED_ACCENT),
    ("Edge\nCases", PURPLE),
    ("Boundary\nTests", BLUE),
]
for i, (txt, col) in enumerate(tc_types):
    add_rounded_box(slide, bx + Inches(i * 1.1), by + Inches(2.5),
                    Inches(1.0), Inches(0.8), txt, col, WHITE, 8)

add_text_label(slide, bx + Inches(0.5), by + Inches(3.6), Inches(3.5), Inches(0.4),
               "3x coverage, 60% less effort", 12, ORANGE, True)


# ═══════════════════════════════════════════════
# SLIDE 10: Stakeholder Communication — Multi-output
# ═══════════════════════════════════════════════
slide = prs.slides[9]
shrink_wide_textbox(slide)

bx = Inches(8.0)
by = Inches(2.0)

add_rounded_box(slide, bx + Inches(0.5), by, Inches(3.5), Inches(0.7),
                "Raw Input / Updates", LIGHT_TEAL, DARK_TEAL, 11)
add_arrow_down(slide, bx + Inches(2.25), by + Inches(0.75), Inches(0.35), TEAL)

add_rounded_box(slide, bx + Inches(0.5), by + Inches(1.2), Inches(3.5), Inches(0.7),
                "AI Communication Engine", TEAL, WHITE, 11)
add_arrow_down(slide, bx + Inches(2.25), by + Inches(1.95), Inches(0.35), TEAL)

comm_outputs = [
    ("Release Notes", GREEN),
    ("Status Reports", BLUE),
    ("Escalation Emails", PURPLE),
    ("Tech Translations", ORANGE),
]
for i, (txt, col) in enumerate(comm_outputs):
    add_rounded_box(slide, bx + Inches(0.2) + Inches((i % 2) * 2.0), by + Inches(2.5 + int(i / 2) * 0.9),
                    Inches(1.9), Inches(0.65), txt, col, WHITE, 9)


# ═══════════════════════════════════════════════
# SLIDE 11: Data Analysis — Question to Answer
# ═══════════════════════════════════════════════
slide = prs.slides[10]
shrink_wide_textbox(slide)

bx = Inches(8.0)
by = Inches(2.0)

add_rounded_box(slide, bx + Inches(0.5), by, Inches(3.5), Inches(0.7),
                "Plain English Question", LIGHT_TEAL, DARK_TEAL, 11)
add_arrow_down(slide, bx + Inches(2.25), by + Inches(0.75), Inches(0.35), TEAL)

add_rounded_box(slide, bx + Inches(0.5), by + Inches(1.2), Inches(3.5), Inches(0.7),
                "AI Data Engine", TEAL, WHITE, 12)
add_arrow_down(slide, bx + Inches(2.25), by + Inches(1.95), Inches(0.35), TEAL)

data_outputs = [
    ("SQL\nQueries", BLUE),
    ("Excel\nFormulas", GREEN),
    ("Pattern\nInsights", PURPLE),
    ("KPI\nSuggestions", ORANGE),
]
for i, (txt, col) in enumerate(data_outputs):
    add_rounded_box(slide, bx + Inches(i * 1.1), by + Inches(2.5),
                    Inches(1.0), Inches(0.8), txt, col, WHITE, 8)

add_text_label(slide, bx + Inches(0.5), by + Inches(3.6), Inches(3.5), Inches(0.4),
               "Hours → Minutes", 13, ORANGE, True)


# ═══════════════════════════════════════════════
# SLIDE 12: Day-to-Day — Time savings circles
# ═══════════════════════════════════════════════
slide = prs.slides[11]
shrink_wide_textbox(slide)

bx = Inches(8.0)
by = Inches(1.8)

add_text_label(slide, bx, by, Inches(4.5), Inches(0.5), "Time Saved Per Task", 13, DARK_TEAL, True)

tasks_circles = [
    ("80%", "Data\nQueries", TEAL),
    ("70%", "Meeting\nNotes", BLUE),
    ("65%", "Jira\nTickets", GREEN),
    ("60%", "Test\nCases", PURPLE),
    ("55%", "PRD\nDrafts", ORANGE),
    ("50%", "Status\nReports", RED_ACCENT),
]

for i, (pct, label, col) in enumerate(tasks_circles):
    row = int(i / 3)
    col_i = i % 3
    cx = bx + Inches(0.3) + Inches(col_i * 1.5)
    cy = by + Inches(0.7) + Inches(row * 2.3)
    size = Inches(1.15)
    add_circle(slide, cx, cy, size, pct, col, WHITE, 16)
    add_text_label(slide, cx - Inches(0.1), cy + size + Inches(0.05), Inches(1.35), Inches(0.4), label, 8, DARK_TEAL, False)


# ═══════════════════════════════════════════════
# SLIDE 13: Before vs After — Comparison arrows
# ═══════════════════════════════════════════════
slide = prs.slides[12]
shrink_wide_textbox(slide)

bx = Inches(8.2)
by = Inches(2.0)

metrics = [
    ("PRD Time", "4-6 hrs", "1-2 hrs"),
    ("Stories/Sprint", "15-20", "30-40"),
    ("Edge Coverage", "60-70%", "85-95%"),
    ("Research", "2-3 days", "3-4 hrs"),
    ("Rework Cycles", "3-4", "1-2"),
]

add_text_label(slide, bx - Inches(0.2), by - Inches(0.5), Inches(1.5), Inches(0.3), "", 9, DARK_TEAL, True)
add_text_label(slide, bx + Inches(1.0), by - Inches(0.5), Inches(1.0), Inches(0.3), "Before", 10, RED_ACCENT, True)
add_text_label(slide, bx + Inches(3.0), by - Inches(0.5), Inches(1.0), Inches(0.3), "After", 10, GREEN, True)

for i, (metric, before, after) in enumerate(metrics):
    y = by + Inches(i * 0.85)
    add_text_label(slide, bx - Inches(0.2), y + Inches(0.05), Inches(1.3), Inches(0.4), metric, 9, DARK_TEAL, True, PP_ALIGN.LEFT)
    add_rounded_box(slide, bx + Inches(1.0), y, Inches(1.0), Inches(0.5), before, RGBColor(0xEF, 0x9A, 0x9A), DARK_TEAL, 10, True)
    add_arrow_right(slide, bx + Inches(2.1), y + Inches(0.25), Inches(0.7), TEAL)
    add_rounded_box(slide, bx + Inches(2.9), y, Inches(1.0), Inches(0.5), after, RGBColor(0xA5, 0xD6, 0xA7), DARK_TEAL, 10, True)


# ═══════════════════════════════════════════════
# SLIDE 14: Tools — Tool cards
# ═══════════════════════════════════════════════
slide = prs.slides[13]
shrink_wide_textbox(slide)

bx = Inches(8.0)
by = Inches(2.0)

tools = [
    ("Claude", "PRDs & Analysis", TEAL),
    ("ChatGPT", "Quick Drafts", BLUE),
    ("Notion AI", "Docs & Wikis", PURPLE),
    ("Copilot", "Tech Specs", GREEN),
    ("Otter.ai", "Transcription", ORANGE),
    ("Gamma", "Presentations", RED_ACCENT),
]

for i, (name, use, col) in enumerate(tools):
    row = int(i / 2)
    col_i = i % 2
    x = bx + Inches(col_i * 2.3)
    y = by + Inches(row * 1.3)
    add_rounded_box(slide, x, y, Inches(2.1), Inches(0.55), name, col, WHITE, 11, True)
    add_text_label(slide, x, y + Inches(0.6), Inches(2.1), Inches(0.3), use, 9, DARK_TEAL, False)


# ═══════════════════════════════════════════════
# SLIDE 15: Best Practices — Numbered circles
# ═══════════════════════════════════════════════
slide = prs.slides[14]
shrink_wide_textbox(slide)

bx = Inches(8.2)
by = Inches(1.8)

practices = [
    ("1", "Review\nAI Output", TEAL),
    ("2", "Give\nContext", BLUE),
    ("3", "Build\nTemplates", PURPLE),
    ("4", "Protect\nData", RED_ACCENT),
    ("5", "Start\nSmall", GREEN),
    ("6", "Share\nLearnings", ORANGE),
]

for i, (num, label, col) in enumerate(practices):
    row = int(i / 3)
    col_i = i % 3
    cx = bx + Inches(0.2) + Inches(col_i * 1.5)
    cy = by + Inches(0.5) + Inches(row * 2.3)
    add_circle(slide, cx, cy, Inches(1.1), num, col, WHITE, 22)
    add_text_label(slide, cx - Inches(0.1), cy + Inches(1.15), Inches(1.3), Inches(0.5), label, 9, DARK_TEAL, True)


# ═══════════════════════════════════════════════
# SLIDE 16: 30-Day Plan — Timeline
# ═══════════════════════════════════════════════
slide = prs.slides[15]
shrink_wide_textbox(slide)

bx = Inches(8.2)
by = Inches(2.0)

weeks = [
    ("Week 1", "Start using AI\nfor 2 tasks", TEAL),
    ("Week 2", "Build prompt\ntemplates", BLUE),
    ("Week 3", "Share & refine\nwith team", PURPLE),
    ("Week 4", "Scale & create\nprompt library", GREEN),
]

for i, (week, desc, col) in enumerate(weeks):
    y = by + Inches(i * 1.15)
    add_circle(slide, bx, y, Inches(0.7), str(i + 1), col, WHITE, 18)
    add_rounded_box(slide, bx + Inches(1.0), y, Inches(3.2), Inches(0.7), f"{week}\n{desc}", col, WHITE, 9)
    if i < len(weeks) - 1:
        # Connecting line
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, bx + Inches(0.3), y + Inches(0.7), Inches(0.1), Inches(0.4)
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = LIGHT_TEAL
        line_shape.line.fill.background()
        line_shape.shadow.inherit = False

add_text_label(slide, bx, by + Inches(4.8), Inches(4.2), Inches(0.5),
               "Goal: Save 5-8 hrs/week", 14, ORANGE, True, PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════
# SLIDE 17: Q&A — Visual emphasis
# ═══════════════════════════════════════════════
slide = prs.slides[16]
shrink_wide_textbox(slide)

bx = Inches(8.5)
by = Inches(2.2)

add_circle(slide, bx + Inches(0.5), by, Inches(2.5), "Q&A", TEAL, WHITE, 32)

pillars = ["Understand\nUsers", "Shape\nProducts", "Drive\nDecisions"]
pillar_colors = [GREEN, BLUE, PURPLE]
for i, (txt, col) in enumerate(zip(pillars, pillar_colors)):
    add_rounded_box(slide, bx + Inches(i * 1.3) - Inches(0.4), by + Inches(3.0),
                    Inches(1.2), Inches(0.8), txt, col, WHITE, 9)

add_text_label(slide, bx - Inches(0.5), by + Inches(4.1), Inches(4.5), Inches(0.4),
               "AI empowers the BA", 13, DARK_TEAL, True)


# ═══════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════
prs.save(OUTPUT)
print(f"✅ PPT with visuals saved to: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
