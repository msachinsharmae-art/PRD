"""
US Payroll Competitor Analysis - PPTX Generator v4
- Keeps original Zimyo title slide (slide 1) and closing slide (slide 8)
- Inner slides use Zimyo template decorative elements (logo, corner shapes)
- Stakeholder-ready language for Indian HRMS company
- Clean grid: left = data, right = visuals
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── SLIDE DIMENSIONS ──
SLIDE_W = 13.333
SLIDE_H = 7.5

# ── GRID (inches) ──
# Zimyo logo sits at X=0.39 Y=0.34, so title starts at X=1.1
TITLE_X = 1.1
TITLE_Y = 0.3
TITLE_W = 10.5
SUB_Y = 0.72      # subtitle line
DIV_Y = 1.0       # divider line
BODY_Y = 1.2      # content starts

# Two-column
LX = 0.7          # left col start
LW = 5.5          # left col width
RX = 6.6          # right col start
RW = 5.9          # right col width

# ── COLORS ──
NAVY = RGBColor(0x03, 0x19, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x64, 0x69, 0x6C)
TXT = RGBColor(0x2D, 0x3A, 0x45)
MUTED = RGBColor(0x6C, 0x75, 0x7D)
PURPLE = RGBColor(0x6C, 0x5C, 0xE7)
TEAL = RGBColor(0x00, 0xB8, 0x94)
CORAL = RGBColor(0xFD, 0x79, 0x72)
BLUE = RGBColor(0x00, 0x84, 0xD6)
ORANGE = RGBColor(0xFF, 0xA5, 0x02)
GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
RED_BG = RGBColor(0xFD, 0xED, 0xED)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)

def e(inches):
    return int(inches * 914400)

# ── SHAPE HELPERS ──

def box(slide, x, y, w, h, fill, line_clr=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, e(x), e(y), e(w), e(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_clr:
        sh.line.color.rgb = line_clr
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    return sh

def rbox(slide, x, y, w, h, fill, line_clr=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, e(x), e(y), e(w), e(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_clr:
        sh.line.color.rgb = line_clr
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.adjustments[0] = 0.06
    return sh

def t(slide, x, y, w, h, text, sz=12, bold=False, color=TXT, align=PP_ALIGN.LEFT):
    bx = slide.shapes.add_textbox(e(x), e(y), e(w), e(h))
    tf = bx.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return bx

def bl(slide, x, y, w, h, items, sz=9, color=TXT, sp=3):
    bx = slide.shapes.add_textbox(e(x), e(y), e(w), e(h))
    tf = bx.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(sp)
        p.space_before = Pt(0)
        r = p.add_run()
        r.text = item
        r.font.size = Pt(sz)
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return bx

def copy_inner_bg(src_slide, tgt_slide):
    """Copy Zimyo decorative elements from template inner slide.
    Copies: logo, corner decorations. Skips: '02' page number, title textbox."""
    skip = {"Google Shape;128;p2", "TextBox 5"}  # 02 number, title
    for shape in src_slide.shapes:
        if shape.name not in skip:
            tgt_slide.shapes._spTree.append(copy.deepcopy(shape._element))

def slide_header(slide, title, subtitle=""):
    """Consistent header on inner slides — works with Zimyo logo at top-left."""
    t(slide, TITLE_X, TITLE_Y, TITLE_W, 0.38, title, sz=20, bold=True, color=NAVY)
    if subtitle:
        t(slide, TITLE_X, SUB_Y, TITLE_W, 0.22, subtitle, sz=9, color=GRAY)
    # Thin accent divider
    box(slide, TITLE_X, DIV_Y, TITLE_W, 0.025, PURPLE)

def section(slide, x, y, w, label, color=NAVY):
    t(slide, x, y, w, 0.22, label, sz=11, bold=True, color=color)

def metric(slide, x, y, w, h, val, label, accent):
    rbox(slide, x, y, w, h, WHITE, line_clr=RGBColor(0xE0, 0xE0, 0xE0))
    box(slide, x + 0.02, y + 0.08, 0.05, h - 0.16, accent)
    t(slide, x + 0.16, y + 0.06, w - 0.25, 0.28, val, sz=15, bold=True, color=accent)
    t(slide, x + 0.16, y + 0.38, w - 0.25, 0.18, label, sz=7, color=MUTED)

def info_box(slide, x, y, w, h, title_text, items, bg, title_clr, txt_clr=TXT):
    rbox(slide, x, y, w, h, bg)
    t(slide, x + 0.12, y + 0.06, w - 0.24, 0.2, title_text, sz=10, bold=True, color=title_clr)
    bl(slide, x + 0.12, y + 0.3, w - 0.24, h - 0.38, items, sz=8, color=txt_clr, sp=2)

def bar_row(slide, x, y, label, pct, accent, bar_w=3.0):
    t(slide, x, y, 1.2, 0.18, label, sz=8, color=TXT)
    bx = x + 1.25
    rbox(slide, bx, y + 0.02, bar_w, 0.13, RGBColor(0xE8, 0xE8, 0xE8))
    rbox(slide, bx, y + 0.02, max(0.06, bar_w * pct / 100), 0.13, accent)
    t(slide, bx + bar_w + 0.08, y, 0.45, 0.18, f"{pct}%", sz=8, bold=True, color=accent)

def add_table(slide, x, y, w, h, data, cws=None):
    nr, nc = len(data), len(data[0])
    ts = slide.shapes.add_table(nr, nc, e(x), e(y), e(w), e(h))
    tbl = ts.table
    if cws:
        for i, cw in enumerate(cws):
            tbl.columns[i].width = e(cw)
    for r in range(nr):
        for c in range(nc):
            cell = tbl.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(data[r][c])
            run.font.name = "Calibri"
            cell.margin_left = e(0.05)
            cell.margin_right = e(0.04)
            cell.margin_top = e(0.02)
            cell.margin_bottom = e(0.02)
            if r == 0:
                run.font.size = Pt(8)
                run.font.bold = True
                run.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            else:
                run.font.size = Pt(7.5)
                run.font.color.rgb = TXT
                if r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_BG


# ═════════════════════════════════════════════
def main():
    prs = Presentation("Zimyo Template_Probation.pptx")
    blank = prs.slide_layouts[2]

    # Save references before deleting
    ref_inner = prs.slides[1]   # inner slide for bg elements
    title_slide = prs.slides[0]
    closing_slide = prs.slides[7]

    # Store XML of title and closing slides
    title_xml = copy.deepcopy(title_slide._element)
    closing_xml = copy.deepcopy(closing_slide._element)
    # Store relationships for title and closing
    title_rels = []
    closing_rels = []
    for rel in title_slide.part.rels.values():
        title_rels.append(rel)
    for rel in closing_slide.part.rels.values():
        closing_rels.append(rel)

    # Store inner slide bg elements (as list of XML copies)
    inner_bg_elements = []
    skip_names = {"Google Shape;128;p2", "TextBox 5"}
    for shape in ref_inner.shapes:
        if shape.name not in skip_names:
            inner_bg_elements.append(copy.deepcopy(shape._element))

    # Store closing slide elements for later recreation
    closing_elements = []
    for shape in closing_slide.shapes:
        closing_elements.append(copy.deepcopy(shape._element))

    # Delete ALL slides except title (index 0)
    # Delete 7 down to 1 in reverse
    for i in range(7, 0, -1):
        rId = prs.slides._sldIdLst[i].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[i])

    # Now: [Title (0)] only
    # Update title slide text
    for shape in prs.slides[0].shapes:
        if shape.name == "Google Shape;98;p1":
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = "US Payroll - Competitor Analysis 2026"

    def new_slide():
        """Add a blank slide and paste inner bg elements."""
        s = prs.slides.add_slide(blank)
        for el in inner_bg_elements:
            s.shapes._spTree.append(copy.deepcopy(el))
        return s

    # ════════════════════════════════════════════
    # SLIDE 2: Market Overview
    # ════════════════════════════════════════════
    s = new_slide()
    slide_header(s, "US Startup Payroll Market - 2026 Overview",
                 "Who is winning in US payroll and what does this mean for us")

    # LEFT
    y = BODY_Y
    section(s, LX, y, LW, "What the Market Looks Like")
    y += 0.3
    bl(s, LX, y, LW, 1.85, [
        "Rippling and Gusto together own 57% of the startup payroll market",
        "Central (by Mercury) and Warp are new entrants - 32% share in 2-3 years",
        "Mercury bundled free payroll with its banking product - fastest growth",
        "Warp offers completely free payroll funded through interest on deposits",
        "Justworks uses PEO (co-employment) model - declining among startups",
        "All top players are based in US, focused on US-first compliance",
    ], sz=9, sp=4)

    y += 2.1
    section(s, LX, y, LW, "What This Means for Zimyo")
    y += 0.3
    bl(s, LX, y, LW, 1.85, [
        "US payroll market is ~$28B and growing 8% yearly in SMB",
        "Nobody is serving Indian companies with US subsidiaries well",
        "Multi-country payroll (India + US) is a gap we can fill",
        "There is room for a new player with competitive pricing",
        "Our existing HRMS gives us a platform advantage to bundle payroll",
    ], sz=9, sp=4)

    # RIGHT - bar chart + metrics
    ry = BODY_Y
    section(s, RX, ry, RW, "Market Share by Provider")
    ry += 0.35
    for name, pct, clr in [
        ("Rippling",  29, PURPLE), ("Gusto",     28, TEAL),
        ("Central",   21, BLUE),   ("Warp",      11, ORANGE),
        ("Deel",       5, CORAL),  ("Justworks",  3, GRAY),
    ]:
        bar_row(s, RX, ry, name, pct, clr, bar_w=3.0)
        ry += 0.28

    ry += 0.25
    section(s, RX, ry, RW, "Key Numbers")
    ry += 0.32
    cw = (RW - 0.2) / 2
    metric(s, RX,        ry, cw, 0.62, "$13.5B", "Rippling Valuation", PURPLE)
    metric(s, RX+cw+0.2, ry, cw, 0.62, "300K+",  "Gusto Customers", TEAL)
    ry += 0.72
    metric(s, RX,        ry, cw, 0.62, "200K+",  "Mercury Accounts", BLUE)
    metric(s, RX+cw+0.2, ry, cw, 0.62, "$12B",   "Deel Valuation", CORAL)
    ry += 0.72
    metric(s, RX,        ry, cw, 0.62, "$0/mo",  "Warp Free Tier", ORANGE)
    metric(s, RX+cw+0.2, ry, cw, 0.62, "35K+",   "Deel Customers", CORAL)

    # ════════════════════════════════════════════
    # SLIDE 3: Rippling
    # ════════════════════════════════════════════
    s = new_slide()
    slide_header(s, "Rippling - 29% Market Share (Largest Player)",
                 "Founded 2016  |  Parker Conrad (ex-Zenefits)  |  San Francisco  |  Valued at $13.5B")

    y = BODY_Y
    section(s, LX, y, LW, "How They Built Their US Payroll", PURPLE)
    y += 0.28
    bl(s, LX, y, LW, 1.75, [
        "Built a single employee database that connects HR, IT and Finance",
        "Payroll engine does automatic tax calculation for all 50 US states",
        "They bought tax filing technology instead of building from scratch",
        "Created a rules engine (RQL) for complex pay rules and deductions",
        "Supports companies with multiple entities and state registrations",
        "Added 600+ integrations with tools like Slack, QuickBooks, Netsuite",
    ], sz=9, sp=3)

    y += 2.0
    section(s, LX, y, LW, "Their US Payroll Covers", PURPLE)
    y += 0.28
    bl(s, LX, y, LW, 1.65, [
        "Federal, state and local tax calculation, filing and payments",
        "Automatic detection when an employee moves to another state",
        "W-2 for employees, 1099 for contractors, new hire reporting",
        "Benefits: health insurance, 401k, HSA, FSA, commuter benefits",
        "Workers compensation and ACA compliance reporting",
        "Global payroll in 50+ countries through EOR and local entities",
    ], sz=9, sp=3)

    y += 1.82
    section(s, LX, y, LW, "Pricing")
    y += 0.25
    t(s, LX, y, LW, 0.18,
      "Starts at $8 per employee per month + platform fee  |  Target: 10 to 2,000 employees", sz=9, color=MUTED)

    # RIGHT
    ry = BODY_Y
    cw = (RW - 0.2) / 2
    metric(s, RX,        ry, cw, 0.6, "30K+", "Companies", PURPLE)
    metric(s, RX+cw+0.2, ry, cw, 0.6, "$1.4B", "Total Funding", PURPLE)
    ry += 0.7
    metric(s, RX,        ry, cw, 0.6, "600+", "Integrations", BLUE)
    metric(s, RX+cw+0.2, ry, cw, 0.6, "50+", "Countries", TEAL)

    ry += 0.85
    info_box(s, RX, ry, RW, 1.35,
        "What Works for Them", [
            "Single platform for HR + IT + Finance saves companies money",
            "Their automation engine is the best in the market right now",
            "Strong compliance across all US states",
            "Very fast at adding new product features",
        ], GREEN_BG, TEAL, TEAL)

    ry += 1.5
    info_box(s, RX, ry, RW, 1.35,
        "Where They Are Weak", [
            "Expensive compared to SMB-focused tools like Gusto",
            "Too complex for small teams with less than 20 people",
            "Customer support gets mixed reviews at scale",
            "Aggressive sales approach puts some buyers off",
        ], RED_BG, CORAL, CORAL)

    ry += 1.55
    rbox(s, RX, ry, RW, 0.55, LIGHT_BG)
    t(s, RX + 0.12, ry + 0.06, RW - 0.24, 0.18,
      "What We Can Learn", sz=9, bold=True, color=NAVY)
    t(s, RX + 0.12, ry + 0.28, RW - 0.24, 0.22,
      "Bundling HR + Payroll on one platform is the winning strategy. Their rules engine approach for compliance is worth studying.",
      sz=8, color=MUTED)

    # ════════════════════════════════════════════
    # SLIDE 4: Gusto
    # ════════════════════════════════════════════
    s = new_slide()
    slide_header(s, "Gusto - 28% Market Share (SMB Favourite)",
                 "Founded 2011 (was ZenPayroll)  |  Josh Reeves CEO  |  San Francisco  |  Valued at $9.5B")

    y = BODY_Y
    section(s, LX, y, LW, "How They Built Their US Payroll", TEAL)
    y += 0.28
    bl(s, LX, y, LW, 1.75, [
        "Started as a payroll-only company (ZenPayroll) and built their tax engine",
        "Registered as tax filing agent in all 50 US states and territories",
        "Built their own benefits brokerage - Gusto is the broker, not a middleman",
        "Acquired Symmetry for their tax calculation technology",
        "Created a channel of 100K+ accounting firms who refer clients",
        "Launched embedded payroll API - other platforms can white-label Gusto",
    ], sz=9, sp=3)

    y += 2.0
    section(s, LX, y, LW, "Their US Payroll Covers", TEAL)
    y += 0.28
    bl(s, LX, y, LW, 1.65, [
        "Unlimited pay runs in any frequency (weekly, biweekly, monthly)",
        "Auto tax filing for federal, state and local across all states",
        "Next-day and 2-day direct deposit options",
        "Health insurance, 401k, HSA, FSA, workers comp (pay-as-you-go)",
        "Contractor payments with automatic 1099-NEC filing",
        "Hiring flow: offer letters, I-9 verification, W-4, state forms",
        "Built-in time tracking and PTO management",
    ], sz=9, sp=3)

    y += 1.95
    section(s, LX, y, LW, "Pricing")
    y += 0.25
    t(s, LX, y, LW, 0.18,
      "Simple: $40/mo + $6/emp  |  Plus: $80/mo + $12/emp  |  Premium: custom  |  Target: 1 to 200 employees",
      sz=9, color=MUTED)

    # RIGHT
    ry = BODY_Y
    cw = (RW - 0.2) / 2
    metric(s, RX,        ry, cw, 0.6, "300K+", "Companies", TEAL)
    metric(s, RX+cw+0.2, ry, cw, 0.6, "6M+",   "Employees Paid", TEAL)
    ry += 0.7
    metric(s, RX,        ry, cw, 0.6, "100K+", "Accounting Partners", BLUE)
    metric(s, RX+cw+0.2, ry, cw, 0.6, "$746M", "Total Funding", ORANGE)

    ry += 0.85
    info_box(s, RX, ry, RW, 1.35,
        "What Works for Them", [
            "Simplest payroll product to set up and use in the market",
            "Being their own broker gives them better margins on benefits",
            "Accountant channel brings them steady customer flow",
            "Embedded API lets other SaaS platforms sell Gusto payroll",
        ], GREEN_BG, TEAL, TEAL)

    ry += 1.5
    info_box(s, RX, ry, RW, 1.35,
        "Where They Are Weak", [
            "Cannot handle companies above 200-300 employees",
            "No IT or device management like Rippling offers",
            "Very limited international payroll support",
            "Adding new features slower than Rippling",
        ], RED_BG, CORAL, CORAL)

    ry += 1.55
    rbox(s, RX, ry, RW, 0.55, LIGHT_BG)
    t(s, RX + 0.12, ry + 0.06, RW - 0.24, 0.18,
      "What We Can Learn", sz=9, bold=True, color=NAVY)
    t(s, RX + 0.12, ry + 0.28, RW - 0.24, 0.22,
      "Their accountant channel is gold - we should build similar CA/CPA partnerships. Embedded payroll API is a smart platform play.",
      sz=8, color=MUTED)

    # ════════════════════════════════════════════
    # SLIDE 5: Central + Warp (side by side)
    # ════════════════════════════════════════════
    s = new_slide()
    slide_header(s, "Fast-Growing Challengers: Central (Mercury) & Warp",
                 "Combined 32% market share in just 2-3 years - they changed the game on pricing")

    card_w = 5.65
    card_h = 5.7
    c1x = 0.7
    c2x = 6.7
    cy = BODY_Y
    pad = 0.18
    iw = card_w - 2 * pad

    # LEFT CARD: Central
    rbox(s, c1x, cy, card_w, card_h, WHITE, line_clr=BLUE)
    ix = c1x + pad
    iy = cy + 0.1

    t(s, ix, iy, iw, 0.28, "Central (Mercury)  |  21% Share", sz=15, bold=True, color=BLUE)
    iy += 0.3
    t(s, ix, iy, iw, 0.18, "Mercury bank added payroll  |  Immad Akhund CEO  |  $1.6B valuation", sz=8, color=MUTED)
    iy += 0.3

    t(s, ix, iy, iw, 0.2, "How They Built It", sz=10, bold=True, color=BLUE)
    iy += 0.24
    bl(s, ix, iy, iw, 1.55, [
        "Mercury started as a banking platform for startups",
        "Added payroll that runs directly from the Mercury bank account",
        "No fund transfer delay - payroll is funded instantly from the bank",
        "200K+ existing banking customers got payroll for free",
        "Tax engine handles federal + all 50 states filing automatically",
        "Also supports contractor payments (1099) and multi-state",
    ], sz=8, sp=3)

    iy += 1.7
    t(s, ix, iy, iw, 0.2, "Why They Grew So Fast", sz=10, bold=True, color=BLUE)
    iy += 0.24
    bl(s, ix, iy, iw, 1.0, [
        "Free payroll bundled with banking - no extra cost for customers",
        "200K existing Mercury users = ready-made distribution",
        "Banking + payroll together is hard for others to copy",
        "Startups see real-time cash flow impact of payroll runs",
    ], sz=8, sp=3)

    iy += 1.15
    rbox(s, ix, iy, iw, 0.3, RED_BG)
    t(s, ix + 0.08, iy + 0.06, iw - 0.16, 0.16,
      "Gaps: Only works if you bank with Mercury  |  Limited HR and benefits features", sz=7.5, color=CORAL)

    # RIGHT CARD: Warp
    rbox(s, c2x, cy, card_w, card_h, WHITE, line_clr=ORANGE)
    ix = c2x + pad
    iy = cy + 0.1

    t(s, ix, iy, iw, 0.28, "Warp  |  11% Share", sz=15, bold=True, color=ORANGE)
    iy += 0.3
    t(s, ix, iy, iw, 0.18, "Free payroll for startups  |  YC and a16z backed  |  ~$28M funding", sz=8, color=MUTED)
    iy += 0.3

    t(s, ix, iy, iw, 0.2, "How They Built It", sz=10, bold=True, color=ORANGE)
    iy += 0.24
    bl(s, ix, iy, iw, 1.55, [
        "Payroll is completely free - they make money from interest on deposits",
        "Customer funds sit in high-yield accounts, Warp keeps the interest",
        "Same model as Robinhood - float-funded revenue",
        "Full tax filing: federal, state, local calculations and payments",
        "Direct deposit, W-2 and 1099 generation, new hire reporting",
        "Multi-state compliance is fully automated",
    ], sz=8, sp=3)

    iy += 1.7
    t(s, ix, iy, iw, 0.2, "Why They Grew So Fast", sz=10, bold=True, color=ORANGE)
    iy += 0.24
    bl(s, ix, iy, iw, 1.0, [
        "$0 per month payroll - hard to compete with free",
        "YC network gave them access to thousands of startups",
        "Very simple product - onboarding takes under 10 minutes",
        "Built specifically for early-stage startup needs",
    ], sz=8, sp=3)

    iy += 1.15
    rbox(s, ix, iy, iw, 0.3, RED_BG)
    t(s, ix + 0.08, iy + 0.06, iw - 0.16, 0.16,
      "Gaps: US only  |  No HR features beyond payroll  |  Small team, limited support", sz=7.5, color=CORAL)

    # ════════════════════════════════════════════
    # SLIDE 6: Deel + Justworks (side by side)
    # ════════════════════════════════════════════
    s = new_slide()
    slide_header(s, "Niche Players: Deel & Justworks",
                 "Deel is global-first (150+ countries)  |  Justworks uses PEO co-employment model")

    # LEFT CARD: Deel
    rbox(s, c1x, cy, card_w, card_h, WHITE, line_clr=CORAL)
    ix = c1x + pad
    iy = cy + 0.1

    t(s, ix, iy, iw, 0.28, "Deel  |  5% Share  (Global-First)", sz=15, bold=True, color=CORAL)
    iy += 0.3
    t(s, ix, iy, iw, 0.18, "Alex Bouaziz CEO  |  $12B val  |  $680M funding  |  35K+ customers", sz=8, color=MUTED)
    iy += 0.3

    t(s, ix, iy, iw, 0.2, "How They Built US Payroll", sz=10, bold=True, color=CORAL)
    iy += 0.24
    bl(s, ix, iy, iw, 1.55, [
        "Started with global EOR (employer of record) in 150+ countries",
        "Added US domestic payroll in 2023 as an expansion",
        "Owns 80+ legal entities worldwide for compliance",
        "Bought PayGroup to strengthen APAC payroll infrastructure",
        "US tax engine covers federal, state and local filing",
        "Launched free HRIS (Deel HR) to compete on platform depth",
    ], sz=8, sp=3)

    iy += 1.7
    t(s, ix, iy, iw, 0.2, "What Makes Them Different", sz=10, bold=True, color=CORAL)
    iy += 0.24
    bl(s, ix, iy, iw, 1.0, [
        "Only player offering global + US payroll in one platform",
        "EOR model lets companies hire anywhere without local entities",
        "Immigration and visa support built into the platform",
        "Pricing: US $49/emp/mo  |  EOR $500+/emp/mo",
    ], sz=8, sp=3)

    iy += 1.15
    rbox(s, ix, iy, iw, 0.3, RED_BG)
    t(s, ix + 0.08, iy + 0.06, iw - 0.16, 0.16,
      "Gaps: US payroll is new, not their main strength  |  Expensive  |  Complex UI", sz=7.5, color=CORAL)

    # RIGHT CARD: Justworks
    rbox(s, c2x, cy, card_w, card_h, WHITE, line_clr=GRAY)
    ix = c2x + pad
    iy = cy + 0.1

    t(s, ix, iy, iw, 0.28, "Justworks  |  3% Share  (PEO)", sz=15, bold=True, color=RGBColor(0x55,0x65,0x75))
    iy += 0.3
    t(s, ix, iy, iw, 0.18, "Isaac Oates CEO  |  NYC  |  $164M funding  |  10K+ companies", sz=8, color=MUTED)
    iy += 0.3

    t(s, ix, iy, iw, 0.2, "How They Built US Payroll", sz=10, bold=True, color=RGBColor(0x55,0x65,0x75))
    iy += 0.24
    bl(s, ix, iy, iw, 1.55, [
        "PEO model: Justworks becomes the co-employer of your staff",
        "This gives small companies access to large-group benefit rates",
        "Fortune 500 level health, dental, vision, 401k at group pricing",
        "Full payroll and tax filing is built into the PEO service",
        "Workers compensation and ACA compliance included",
        "12+ years in market - very mature and tested product",
    ], sz=8, sp=3)

    iy += 1.7
    t(s, ix, iy, iw, 0.2, "What Makes Them Different", sz=10, bold=True, color=RGBColor(0x55,0x65,0x75))
    iy += 0.24
    bl(s, ix, iy, iw, 1.0, [
        "PEO gives small companies enterprise-grade benefits",
        "Flat pricing: Basic $59/emp  |  Plus $99/emp with benefits",
        "Very transparent pricing with no hidden fees",
        "Strong compliance and HR support included",
    ], sz=8, sp=3)

    iy += 1.15
    rbox(s, ix, iy, iw, 0.3, RED_BG)
    t(s, ix + 0.08, iy + 0.06, iw - 0.16, 0.16,
      "Gaps: PEO means less control for employer  |  3% shows declining model  |  No global payroll", sz=7.5, color=CORAL)

    # ════════════════════════════════════════════
    # SLIDE 7: Feature Comparison
    # ════════════════════════════════════════════
    s = new_slide()
    slide_header(s, "Head-to-Head Feature Comparison",
                 "How each provider stacks up on core US payroll capabilities")

    add_table(s, 0.7, BODY_Y, 11.9, 3.85, [
        ["US Payroll Capability", "Rippling", "Gusto", "Central", "Warp", "Deel", "Justworks"],
        ["Multi-State Tax Filing", "All 50 states", "All 50 states", "All 50 states", "All 50 states", "All 50 states", "Via PEO"],
        ["Tax Engine", "Acquired + built", "Acquired Symmetry", "Built in-house", "Built in-house", "Acquired + built", "PEO partner"],
        ["Direct Deposit Speed", "Same-day option", "Next-day / 2-day", "Instant (Mercury)", "Standard ACH", "Standard ACH", "Standard ACH"],
        ["Benefits Admin", "Full suite", "Own brokerage", "Through partners", "Through partners", "Through partners", "PEO group rates"],
        ["Contractor (1099)", "Yes", "Yes + international", "Yes", "Yes", "150+ countries", "Limited"],
        ["Global Payroll", "50+ countries", "20+ (contractors)", "No", "No", "150+ countries", "No"],
        ["Compliance Depth", "Advanced + auto", "Strong + AI alerts", "Basic", "Basic", "Global + US", "PEO handles all"],
        ["API / Embedded", "600+ integrations", "Embedded API", "API available", "API available", "API available", "Limited"],
        ["Pricing", "$8/emp/mo", "$6-12/emp/mo", "Free (bundled)", "Free", "$49/emp/mo", "$59-99/emp/mo"],
    ], cws=[1.9, 1.67, 1.67, 1.67, 1.67, 1.67, 1.67])

    # Takeaway box
    iy = 5.3
    rbox(s, 0.7, iy, 11.9, 1.5, NAVY)
    t(s, 0.9, iy + 0.1, 5.5, 0.2, "What This Comparison Tells Us", sz=11, bold=True, color=WHITE)
    bl(s, 0.9, iy + 0.38, 5.5, 1.0, [
        "Tax filing in all 50 states is a basic requirement - everyone has it",
        "Real competition is on benefits, global reach and pricing",
        "Direct deposit speed is becoming a selling point",
    ], sz=8, color=RGBColor(0xBB, 0xBB, 0xBB), sp=3)
    bl(s, 6.6, iy + 0.38, 5.5, 1.0, [
        "Embedded payroll API is the next platform play",
        "Nobody excels at everything - there are clear gaps to exploit",
        "India-US multi-country payroll is completely unserved",
    ], sz=8, color=RGBColor(0xBB, 0xBB, 0xBB), sp=3)

    # ════════════════════════════════════════════
    # SLIDE 8: What to Build
    # ════════════════════════════════════════════
    s = new_slide()
    slide_header(s, "What We Need to Build for US Payroll",
                 "Key infrastructure, compliance and partnerships based on how competitors did it")

    # LEFT - Must build items
    y = BODY_Y
    section(s, LX, y, LW, "Must-Build Components (Priority Order)")
    y += 0.32

    items = [
        ("1. Tax Calculation Engine", CORAL,
         "Calculate federal, state and local taxes for all 50 states. Handle supplemental wages, pre-tax deductions, garnishments. Option: license from Symmetry or Vertex instead of building."),
        ("2. Tax Filing and Payment", CORAL,
         "Register as tax agent in all 50 states (takes 3-6 months). File quarterly 941, annual W-2/1099. Pay taxes to IRS and state agencies on time. We are liable for any errors."),
        ("3. Direct Deposit (ACH)", ORANGE,
         "Partner with Column, Moov or Modern Treasury for ACH processing. Need NACHA certification. Handle payment returns and reversals. Same-day ACH is becoming expected."),
        ("4. Multi-State Compliance", ORANGE,
         "State reciprocity rules, local tax rules (NYC, SF, Philadelphia etc). Auto-detect when employee moves states. New hire reporting requirements vary by state."),
        ("5. Benefits Integration", BLUE,
         "Health insurance, 401k, HSA/FSA, workers comp. Partner with existing brokers first (SimplyInsured, Vestwell). Benefits deductions feed directly into payroll calculations."),
        ("6. Year-End Processing", BLUE,
         "W-2 for employees, 1099-NEC for contractors, ACA 1095-C reports. Corrections handling (W-2c). IRS has strict deadlines with zero tolerance."),
    ]
    for title_text, clr, desc in items:
        box(s, LX, y, 0.05, 0.56, clr)
        t(s, LX + 0.14, y, LW - 0.14, 0.18, title_text, sz=9, bold=True, color=NAVY)
        t(s, LX + 0.14, y + 0.2, LW - 0.14, 0.36, desc, sz=7.5, color=MUTED)
        y += 0.66

    # RIGHT - Priority cards
    ry = BODY_Y
    section(s, RX, ry, RW, "Priority and Complexity")
    ry += 0.32

    priorities = [
        ("Tax Engine",       "MUST HAVE", "License Symmetry/Vertex or build. Every rupee must be accurate.", CORAL),
        ("Tax Filing",       "MUST HAVE", "Register in all 50 states. Start 6 months before launch date.", CORAL),
        ("ACH Payments",     "MUST HAVE", "Partner with Column or Moov. Do not build banking infra ourselves.", CORAL),
        ("State Compliance", "HIGH",      "Garnishments, minimum wage, overtime, sick leave rules per state.", ORANGE),
        ("Benefits",         "HIGH",      "Partner first with existing brokers. Build own brokerage later.", ORANGE),
        ("Onboarding",       "HIGH",      "I-9, W-4, state tax forms, offer letters. Required for good UX.", BLUE),
        ("Contractor 1099",  "MEDIUM",    "Simpler than W-2 payroll. 1099-NEC filing and payments.", BLUE),
        ("Embedded API",     "PHASE 2",   "Let other platforms sell our payroll. Good revenue channel.", TEAL),
        ("Global Payroll",   "PHASE 3",   "Use EOR partnerships for international. Do not build on day one.", GRAY),
    ]

    for name, pri, note, clr in priorities:
        rbox(s, RX, ry, RW, 0.52, WHITE, line_clr=clr)
        t(s, RX + 0.1, ry + 0.04, 1.7, 0.16, name, sz=8, bold=True, color=NAVY)
        # Badge
        bx = RX + RW - 1.05
        rbox(s, bx, ry + 0.04, 0.9, 0.18, clr)
        t(s, bx + 0.04, ry + 0.04, 0.82, 0.18, pri, sz=6.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        t(s, RX + 0.1, ry + 0.25, RW - 0.2, 0.22, note, sz=7, color=MUTED)
        ry += 0.58

    # ════════════════════════════════════════════
    # SLIDE 9: Our Strategy
    # ════════════════════════════════════════════
    s = new_slide()
    slide_header(s, "Our Plan to Enter US Payroll",
                 "Phased approach using our India strengths and existing HRMS platform")

    # LEFT - 3 phases
    y = BODY_Y
    section(s, LX, y, LW, "Phased Roadmap")
    y += 0.32

    # Phase 1
    rbox(s, LX, y, LW, 1.45, WHITE, line_clr=PURPLE)
    t(s, LX + 0.15, y + 0.08, LW - 0.3, 0.22,
      "Phase 1: India-US Corridor  (0-12 months)", sz=11, bold=True, color=PURPLE)
    bl(s, LX + 0.15, y + 0.35, LW - 0.3, 1.0, [
        "Target our existing customers who have US subsidiaries",
        "Build core US payroll: tax calculation, ACH, W-2/1099",
        "Our advantage: dual-country payroll nobody else does well",
        "Price at $10-15 per employee (30-50% below Rippling/Gusto)",
    ], sz=8, sp=3)
    y += 1.6

    # Phase 2
    rbox(s, LX, y, LW, 1.45, WHITE, line_clr=TEAL)
    t(s, LX + 0.15, y + 0.08, LW - 0.3, 0.22,
      "Phase 2: US Startup Market  (12-24 months)", sz=11, bold=True, color=TEAL)
    bl(s, LX + 0.15, y + 0.35, LW - 0.3, 1.0, [
        "Open to general US startups and SMBs (1 to 200 employees)",
        "Add benefits administration, time tracking, hiring tools",
        "Build accountant partnerships and embedded payroll API",
        "Compete on price + multi-country payroll (India, US, UAE, SEA)",
    ], sz=8, sp=3)
    y += 1.6

    # Phase 3
    rbox(s, LX, y, LW, 1.45, WHITE, line_clr=ORANGE)
    t(s, LX + 0.15, y + 0.08, LW - 0.3, 0.22,
      "Phase 3: Mid-Market Push  (24-36 months)", sz=11, bold=True, color=ORANGE)
    bl(s, LX + 0.15, y + 0.35, LW - 0.3, 1.0, [
        "Target mid-market companies (200-1000 employees)",
        "Add advanced compliance, custom reporting, EOR partnerships",
        "Position as the only platform doing India + US + SEA natively",
        "Fill the gap between Gusto (tops at 200) and ADP/Workday",
    ], sz=8, sp=3)

    # RIGHT - Our edges + target
    ry = BODY_Y
    section(s, RX, ry, RW, "Why We Can Win")
    ry += 0.32

    edges = [
        ("India-US Corridor",
         "10,000+ Indian companies have US operations. No payroll provider serves this well today.",
         PURPLE),
        ("Cost Advantage",
         "Our India-based engineering gives us 60-70% lower development cost. We pass this saving to customers.",
         TEAL),
        ("Existing HRMS Platform",
         "We already have HR, Leave, Attendance, Performance modules. Adding payroll makes us all-in-one.",
         BLUE),
        ("Multi-Country Experience",
         "Our India payroll compliance knowledge (PF, ESI, TDS) gives us a framework for US compliance.",
         ORANGE),
        ("Speed to Market",
         "License tax engine (Symmetry), partner for benefits. We do not need to build everything ourselves.",
         CORAL),
    ]

    for title_text, desc, clr in edges:
        rbox(s, RX, ry, RW, 0.78, WHITE, line_clr=clr)
        box(s, RX + 0.02, ry + 0.1, 0.04, 0.58, clr)
        t(s, RX + 0.16, ry + 0.06, RW - 0.28, 0.2, title_text, sz=9, bold=True, color=clr)
        t(s, RX + 0.16, ry + 0.3, RW - 0.28, 0.42, desc, sz=8, color=MUTED)
        ry += 0.87

    # Target box
    ry += 0.08
    rbox(s, RX, ry, RW, 0.8, NAVY)
    t(s, RX + 0.15, ry + 0.08, RW - 0.3, 0.2,
      "24-Month Target", sz=12, bold=True, color=WHITE)
    t(s, RX + 0.15, ry + 0.32, RW - 0.3, 0.18,
      "500+ companies  |  10,000+ employees  |  $5M ARR", sz=10, color=RGBColor(0xBB, 0xBB, 0xBB))
    t(s, RX + 0.15, ry + 0.54, RW - 0.3, 0.18,
      "First mover in India-US payroll corridor", sz=9, bold=True, color=PURPLE)

    # ── Recreate closing slide as last slide ──
    closing = prs.slides.add_slide(blank)
    for el in closing_elements:
        closing.shapes._spTree.append(copy.deepcopy(el))

    # Save
    out = "US_Payroll_Competitor_Analysis_v2.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")
    for i, sl in enumerate(prs.slides):
        txts = []
        for sh in sl.shapes:
            if hasattr(sh, 'text') and sh.text:
                tx = sh.text.strip()[:55].replace('\n',' ').replace('\x0b',' ')
                if tx and len(tx) > 3:
                    txts.append(tx)
        print(f"  {i+1}. {txts[0] if txts else '[visual]'}")

if __name__ == "__main__":
    main()
